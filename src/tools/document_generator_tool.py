"""Document and Presentation Generator Tool"""
import os
import json
import logging
from pathlib import Path
from typing import Dict, Any, List
from datetime import datetime
from docx import Document
from docx.shared import Inches, Pt, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from pptx import Presentation
from pptx.util import Inches as PptxInches, Pt as PptxPt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor as PptxRGBColor
import requests

logger = logging.getLogger(__name__)

class DocumentGeneratorTool:
    """Generate professional documents and presentations with AI research"""
    
    def __init__(self, ai_client=None):
        self.ai_client = ai_client
        self.documents_dir = Path("documents")
        self.documents_dir.mkdir(exist_ok=True)
    
    async def create_document_and_presentation(
        self,
        request: str,
        project_name: str = None,
        output_folder: str = None
    ) -> Dict[str, Any]:
        """Complete workflow: Research → Word Doc → PowerPoint → Save Locally"""
        try:
            logger.info(f"Starting document generation: {request}")
            
            # Generate project name if not provided
            if not project_name:
                project_name = self._sanitize_filename(request[:50])
            
            # Setup project folder
            if output_folder:
                project_path = Path(output_folder) / project_name
            else:
                project_path = self.documents_dir / project_name
            
            if project_path.exists():
                timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
                project_path = project_path.parent / f"{project_name}_{timestamp}"
            
            project_path.mkdir(parents=True, exist_ok=True)
            logger.info(f"✅ Created project directory: {project_path}")
            
            # Step 1: Research and gather information
            research_data = await self._research_topic(request)
            logger.info("✅ Completed research")
            
            # Save research data
            research_path = project_path / "research_data.json"
            with open(research_path, "w", encoding="utf-8") as f:
                json.dump(research_data, f, indent=2, ensure_ascii=False)
            
            # Step 2: Generate structured content
            structured_content = await self._structure_content(request, research_data)
            logger.info("✅ Structured content")
            
            # Step 3: Create Word document
            docx_path = project_path / f"{self._sanitize_filename(request[:30])}.docx"
            await self._create_word_document(docx_path, request, structured_content)
            logger.info(f"✅ Created Word document: {docx_path.name}")
            
            # Step 4: Create PowerPoint presentation
            pptx_path = project_path / f"{self._sanitize_filename(request[:30])}_presentation.pptx"
            await self._create_powerpoint(pptx_path, request, structured_content)
            logger.info(f"✅ Created PowerPoint: {pptx_path.name}")
            
            # Step 5: Save AI prompt history
            prompt_history_path = project_path / "ai_prompt_history.txt"
            with open(prompt_history_path, "w", encoding="utf-8") as f:
                f.write(f"Original Request: {request}\n")
                f.write(f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}\n\n")
                f.write("Research Queries:\n")
                for query in research_data.get("queries", []):
                    f.write(f"  - {query}\n")
            
            # Step 6: Optional PDF export
            pdf_folder = project_path / "exported_pdfs"
            pdf_folder.mkdir(exist_ok=True)
            
            # Summary
            return {
                "success": True,
                "message": f"Documents created successfully!",
                "data": {
                    "project_name": project_name,
                    "local_path": str(project_path),
                    "word_document": str(docx_path),
                    "powerpoint": str(pptx_path),
                    "research_data": str(research_path),
                    "total_slides": structured_content.get("slide_count", 0),
                    "word_count": structured_content.get("word_count", 0)
                }
            }
            
        except Exception as e:
            logger.error(f"Failed to create documents: {e}", exc_info=True)
            return {"success": False, "error": str(e)}
    
    async def _research_topic(self, request: str) -> Dict:
        """Research the topic using AI and web search"""
        try:
            # If AI client available with web search, use it
            research_queries = [
                f"{request} overview",
                f"{request} key features",
                f"{request} market analysis",
                f"{request} best practices"
            ]
            
            research_results = {
                "topic": request,
                "queries": research_queries,
                "findings": [],
                "key_points": [],
                "statistics": [],
                "generated_at": datetime.now().isoformat()
            }
            
            # Simulated research data (in production, use actual AI/web search)
            research_results["findings"] = [
                {
                    "title": "Market Overview",
                    "content": f"The {request} market shows significant growth potential with increasing demand."
                },
                {
                    "title": "Key Features",
                    "content": "Essential features include scalability, user-friendly interface, and robust security."
                },
                {
                    "title": "Target Audience",
                    "content": "Primary users include businesses and individual consumers seeking innovative solutions."
                }
            ]
            
            research_results["key_points"] = [
                f"Growing market demand for {request}",
                "Focus on user experience and accessibility",
                "Integration with existing systems",
                "Cost-effective implementation",
                "Scalable architecture"
            ]
            
            research_results["statistics"] = [
                {"metric": "Market Size", "value": "$X Billion"},
                {"metric": "Growth Rate", "value": "XX% annually"},
                {"metric": "Adoption Rate", "value": "XX%"}
            ]
            
            return research_results
            
        except Exception as e:
            logger.error(f"Research failed: {e}")
            return {
                "topic": request,
                "findings": [],
                "key_points": [],
                "error": str(e)
            }
    
    async def _structure_content(self, request: str, research_data: Dict) -> Dict:
        """Structure content into sections and slides"""
        
        # Document structure
        sections = [
            {
                "title": "Executive Summary",
                "content": f"This document provides a comprehensive overview of {request}, including market analysis, key features, implementation strategy, and expected outcomes."
            },
            {
                "title": "Introduction",
                "content": f"In today's rapidly evolving landscape, {request} has emerged as a critical solution. This proposal outlines the key aspects and strategic approach."
            },
            {
                "title": "Market Analysis",
                "content": "Current market trends indicate strong demand and growth potential.",
                "bullets": research_data.get("key_points", [])
            },
            {
                "title": "Key Features",
                "content": "The solution includes several essential features:",
                "bullets": [
                    "User-friendly interface",
                    "Scalable architecture",
                    "Robust security measures",
                    "Integration capabilities",
                    "Real-time analytics"
                ]
            },
            {
                "title": "Implementation Plan",
                "content": "Phased approach to ensure smooth deployment:",
                "bullets": [
                    "Phase 1: Planning and Design (2-3 weeks)",
                    "Phase 2: Development (4-6 weeks)",
                    "Phase 3: Testing and QA (2 weeks)",
                    "Phase 4: Deployment and Training (1-2 weeks)"
                ]
            },
            {
                "title": "Budget and Timeline",
                "content": "Resource allocation and project timeline details.",
                "bullets": [
                    "Total Budget: $XXX,XXX",
                    "Timeline: 3-4 months",
                    "Team Size: X members",
                    "ROI Expected: XX% within 12 months"
                ]
            },
            {
                "title": "Conclusion",
                "content": f"The proposed {request} solution offers significant value and aligns with strategic objectives. We recommend proceeding with implementation as outlined."
            }
        ]
        
        # Calculate word count
        word_count = sum(
            len(section.get("content", "").split()) +
            sum(len(bullet.split()) for bullet in section.get("bullets", []))
            for section in sections
        )
        
        return {
            "title": request,
            "sections": sections,
            "word_count": word_count,
            "slide_count": len(sections) + 1  # +1 for title slide
        }
    
    async def _create_word_document(
        self,
        file_path: Path,
        request: str,
        content: Dict
    ):
        """Create professionally formatted Word document"""
        
        doc = Document()
        
        # Title page
        title = doc.add_heading(request.title(), 0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        
        doc.add_paragraph()
        
        subtitle = doc.add_paragraph("Project Proposal")
        subtitle.alignment = WD_ALIGN_PARAGRAPH.CENTER
        subtitle.runs[0].font.size = Pt(16)
        subtitle.runs[0].font.color.rgb = RGBColor(100, 100, 100)
        
        doc.add_paragraph()
        doc.add_paragraph()
        
        date_para = doc.add_paragraph(f"Generated: {datetime.now().strftime('%B %d, %Y')}")
        date_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
        date_para.runs[0].font.size = Pt(12)
        
        doc.add_page_break()
        
        # Table of Contents
        doc.add_heading("Table of Contents", 1)
        for i, section in enumerate(content["sections"], 1):
            toc_item = doc.add_paragraph(f"{i}. {section['title']}")
            toc_item.runs[0].font.size = Pt(12)
        
        doc.add_page_break()
        
        # Content sections
        for section in content["sections"]:
            # Section heading
            doc.add_heading(section["title"], 1)
            
            # Section content
            if "content" in section:
                para = doc.add_paragraph(section["content"])
                para.runs[0].font.size = Pt(11)
            
            # Bullets
            if "bullets" in section:
                doc.add_paragraph()
                for bullet in section["bullets"]:
                    bullet_para = doc.add_paragraph(bullet, style='List Bullet')
                    bullet_para.runs[0].font.size = Pt(11)
            
            doc.add_paragraph()
        
        # Save document
        doc.save(str(file_path))
    
    async def _create_powerpoint(
        self,
        file_path: Path,
        request: str,
        content: Dict
    ):
        """Create professional PowerPoint presentation"""
        
        prs = Presentation()
        prs.slide_width = PptxInches(10)
        prs.slide_height = PptxInches(7.5)
        
        # Title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = request.title()
        subtitle.text = f"Project Proposal\n{datetime.now().strftime('%B %Y')}"
        
        # Content slides
        for section in content["sections"]:
            # Content slide layout
            bullet_slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(bullet_slide_layout)
            
            # Title
            title = slide.shapes.title
            title.text = section["title"]
            
            # Content
            text_frame = slide.placeholders[1].text_frame
            text_frame.clear()
            
            # Add paragraph content
            if "content" in section and section["content"]:
                p = text_frame.paragraphs[0]
                p.text = section["content"]
                p.font.size = PptxPt(16)
                p.level = 0
            
            # Add bullets
            if "bullets" in section:
                for bullet in section["bullets"]:
                    p = text_frame.add_paragraph()
                    p.text = bullet
                    p.font.size = PptxPt(14)
                    p.level = 1
        
        # Final slide - Thank You
        final_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(final_slide_layout)
        title = slide.shapes.title
        title.text = "Thank You"
        
        subtitle = slide.placeholders[1]
        subtitle.text = "Questions?"
        
        # Save presentation
        prs.save(str(file_path))
    
    def _sanitize_filename(self, filename: str) -> str:
        """Sanitize filename for safe file system use"""
        import re
        # Remove invalid characters
        filename = re.sub(r'[<>:"/\\|?*]', '', filename)
        # Replace spaces with underscores
        filename = filename.replace(' ', '_')
        # Remove multiple underscores
        filename = re.sub(r'_+', '_', filename)
        return filename.strip('_')